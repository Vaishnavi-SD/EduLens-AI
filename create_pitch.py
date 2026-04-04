from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

GOOGLE_BLUE = RGBColor(0x1a, 0x73, 0xe8)
GOOGLE_GREEN = RGBColor(0x0f, 0x9d, 0x58)
GOOGLE_RED = RGBColor(0xdb, 0x44, 0x37)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x20, 0x21, 0x24)
LIGHT_GRAY = RGBColor(0xF8, 0xF9, 0xFA)

def add_slide(prs, layout_idx=6):
    layout = prs.slide_layouts[layout_idx]
    return prs.slides.add_slide(layout)

def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text(slide, text, left, top, width, height,
             size=24, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, italic=False):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox

def add_rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(
        1, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

# ============================================================
# SLIDE 1 — TITLE
# ============================================================
slide = add_slide(prs)
set_bg(slide, DARK)
add_rect(slide, 0, 0, 13.33, 0.08, GOOGLE_BLUE)
add_rect(slide, 0, 7.42, 13.33, 0.08, GOOGLE_BLUE)

add_text(slide, "EduLens AI", 1, 1.2, 11, 1.5,
         size=54, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(slide, "Intelligent Student Dropout Prevention",
         1, 2.7, 11, 0.8, size=24, color=GOOGLE_BLUE,
         align=PP_ALIGN.CENTER)
add_text(slide, "for Google Classroom",
         1, 3.3, 11, 0.8, size=24, color=GOOGLE_BLUE,
         align=PP_ALIGN.CENTER)
add_text(slide, "Built by Vaishnavi Singasani  |  AI Engineer & Data Analyst",
         1, 4.5, 11, 0.6, size=16, color=RGBColor(0xBB, 0xBB, 0xBB),
         align=PP_ALIGN.CENTER)
add_text(slide, "Published CNN Researcher  |  10M-User App Developer",
         1, 5.0, 11, 0.6, size=16, color=RGBColor(0xBB, 0xBB, 0xBB),
         align=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 2 — THE PROBLEM
# ============================================================
slide = add_slide(prs)
set_bg(slide, DARK)
add_rect(slide, 0, 0, 13.33, 0.08, GOOGLE_RED)

add_text(slide, "The Problem", 0.5, 0.3, 12, 0.8,
         size=36, bold=True, color=WHITE)
add_text(slide, "Google Classroom has 170M+ users.",
         0.8, 1.3, 11, 0.5, size=22, color=WHITE)
add_text(slide, "It knows who logged in.", 0.8, 1.8,
         11, 0.5, size=22, color=WHITE)
add_text(slide, "It does NOT know who is about to quit.",
         0.8, 2.3, 11, 0.5, size=22, bold=True,
         color=RGBColor(0xdb, 0x44, 0x37))

add_text(slide, "1 in 5 online students drops out before completing",
         1.2, 3.2, 10, 0.45, size=18,
         color=RGBColor(0xBB, 0xBB, 0xBB))
add_text(slide, "Teachers cannot monitor 40+ students individually",
         1.2, 3.65, 10, 0.45, size=18,
         color=RGBColor(0xBB, 0xBB, 0xBB))
add_text(slide, "By the time dropout is noticed — it is too late",
         1.2, 4.1, 10, 0.45, size=18,
         color=RGBColor(0xBB, 0xBB, 0xBB))

add_rect(slide, 0.5, 5.2, 12.3, 0.8, RGBColor(0x2a, 0x2a, 0x2a))
add_text(slide,
         "Google Classroom is missing an intelligence layer.",
         0.7, 5.3, 12, 0.6, size=20, bold=True, italic=True,
         color=GOOGLE_BLUE, align=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 3 — THE SOLUTION
# ============================================================
slide = add_slide(prs)
set_bg(slide, DARK)
add_rect(slide, 0, 0, 13.33, 0.08, GOOGLE_GREEN)

add_text(slide, "The Solution", 0.5, 0.3, 12, 0.8,
         size=36, bold=True, color=WHITE)
add_text(slide, "EduLens AI — Real-Time Student Intelligence Platform",
         0.5, 1.1, 12, 0.6, size=22, bold=True, color=GOOGLE_GREEN)

items = [
    ("1", "Cognitive State Detection",
     "CNN model detects if student is engaged, confused, bored or frustrated"),
    ("2", "Dropout Prediction",
     "XGBoost model predicts dropout 2 weeks before it happens — 80%+ accuracy"),
    ("3", "Personalized Recovery Path",
     "AI engine auto-generates a learning plan tailored to each student's style"),
]
for i, (num, title, desc) in enumerate(items):
    y = 2.0 + i * 1.5
    add_rect(slide, 0.5, y, 0.7, 0.7, GOOGLE_BLUE)
    add_text(slide, num, 0.5, y, 0.7, 0.7,
             size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, title, 1.4, y, 11, 0.45,
             size=20, bold=True, color=WHITE)
    add_text(slide, desc, 1.4, y + 0.42, 11, 0.5,
             size=15, color=RGBColor(0xBB, 0xBB, 0xBB))

# ============================================================
# SLIDE 4 — LIVE DEMO
# ============================================================
slide = add_slide(prs)
set_bg(slide, DARK)
add_rect(slide, 0, 0, 13.33, 0.08, GOOGLE_BLUE)

add_text(slide, "Live Demo", 0.5, 0.3, 12, 0.8,
         size=36, bold=True, color=WHITE)
add_rect(slide, 0.5, 1.2, 5.8, 3.5, RGBColor(0x2a, 0x2a, 0x2a))
add_text(slide, "Streamlit Live Prediction App",
         0.5, 2.5, 5.8, 0.6, size=16,
         color=GOOGLE_BLUE, align=PP_ALIGN.CENTER)
add_rect(slide, 6.8, 1.2, 5.8, 3.5, RGBColor(0x2a, 0x2a, 0x2a))
add_text(slide, "Tableau Educator Dashboard",
         6.8, 2.5, 5.8, 0.6, size=16,
         color=GOOGLE_GREEN, align=PP_ALIGN.CENTER)
add_text(slide,
         "Real predictions. Real interventions. Built and running today.",
         0.5, 5.2, 12.3, 0.6, size=20, bold=True, italic=True,
         color=WHITE, align=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 5 — HOW IT WORKS
# ============================================================
slide = add_slide(prs)
set_bg(slide, DARK)
add_rect(slide, 0, 0, 13.33, 0.08, GOOGLE_BLUE)

add_text(slide, "How It Works", 0.5, 0.3, 12, 0.8,
         size=36, bold=True, color=WHITE)

pipeline = [
    ("Student\nSessions", GOOGLE_BLUE),
    ("SQL\nDatabase", RGBColor(0x53, 0x4A, 0xB7)),
    ("XGBoost\nModel", GOOGLE_RED),
    ("Path\nEngine", GOOGLE_GREEN),
    ("Tableau\nDashboard", RGBColor(0xBA, 0x75, 0x17)),
]
for i, (label, color) in enumerate(pipeline):
    x = 0.5 + i * 2.5
    add_rect(slide, x, 1.8, 2.0, 1.2, color)
    add_text(slide, label, x, 1.9, 2.0, 1.0,
             size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    if i < len(pipeline) - 1:
        add_text(slide, "→", x + 2.0, 2.1, 0.5, 0.8,
                 size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_text(slide, "Tech Stack: Python  •  SQL  •  XGBoost  •  Streamlit  •  Tableau  •  SQLite",
         0.5, 3.6, 12, 0.6, size=16,
         color=RGBColor(0xBB, 0xBB, 0xBB), align=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 6 — WHY GOOGLE
# ============================================================
slide = add_slide(prs)
set_bg(slide, DARK)
add_rect(slide, 0, 0, 13.33, 0.08, GOOGLE_BLUE)

add_text(slide, "Why Google", 0.5, 0.3, 12, 0.8,
         size=36, bold=True, color=WHITE)

points = [
    (GOOGLE_BLUE,  "Google Classroom",
     "170M+ users across 180 countries — EduLens plugs in directly"),
    (GOOGLE_GREEN, "Google Meet",
     "Perfect video input for CNN cognitive state detection"),
    (GOOGLE_RED,   "Gemini API",
     "Powers dynamic, intelligent learning path generation"),
    (GOOGLE_BLUE,  "Google Cloud",
     "Scales EduLens to every school on the planet"),
]
for i, (color, title, desc) in enumerate(points):
    y = 1.4 + i * 1.3
    add_rect(slide, 0.5, y, 0.15, 0.7, color)
    add_text(slide, title, 0.9, y, 4, 0.4,
             size=18, bold=True, color=WHITE)
    add_text(slide, desc, 0.9, y + 0.38, 11, 0.45,
             size=15, color=RGBColor(0xBB, 0xBB, 0xBB))

add_rect(slide, 0.5, 6.5, 12.3, 0.7, RGBColor(0x1a, 0x73, 0xe8))
add_text(slide,
         "EduLens is what Google Classroom looks like with Gemini inside.",
         0.7, 6.55, 12, 0.55, size=18, bold=True, italic=True,
         color=WHITE, align=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 7 — MARKET OPPORTUNITY
# ============================================================
slide = add_slide(prs)
set_bg(slide, DARK)
add_rect(slide, 0, 0, 13.33, 0.08, GOOGLE_GREEN)

add_text(slide, "Market Opportunity", 0.5, 0.3, 12, 0.8,
         size=36, bold=True, color=WHITE)

stats = [
    ("$340B", "Global EdTech market by 2025", GOOGLE_BLUE),
    ("$30B",  "Lost annually to online dropout", GOOGLE_RED),
    ("170M+", "Google Classroom active users", GOOGLE_GREEN),
    ("20%",   "Google for Education YoY growth", RGBColor(0xBA, 0x75, 0x17)),
]
for i, (num, label, color) in enumerate(stats):
    x = 0.4 + i * 3.1
    add_rect(slide, x, 1.5, 2.8, 2.0, RGBColor(0x2a, 0x2a, 0x2a))
    add_text(slide, num, x, 1.7, 2.8, 0.9,
             size=36, bold=True, color=color, align=PP_ALIGN.CENTER)
    add_text(slide, label, x, 2.7, 2.8, 0.7,
             size=13, color=RGBColor(0xBB, 0xBB, 0xBB),
             align=PP_ALIGN.CENTER)

add_text(slide,
         "EduLens addresses the #1 unsolved problem in education technology.",
         0.5, 4.2, 12.3, 0.6, size=20, bold=True,
         color=WHITE, align=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 8 — CREDENTIALS
# ============================================================
slide = add_slide(prs)
set_bg(slide, DARK)
add_rect(slide, 0, 0, 13.33, 0.08, GOOGLE_BLUE)

add_text(slide, "Why I Built This", 0.5, 0.3, 12, 0.8,
         size=36, bold=True, color=WHITE)

creds = [
    "B.Tech CSE Data Science — Institute of Aeronautical Engineering, Hyderabad",
    "Published CNN Research — JISEM Journal (84.73% accuracy, real deployment)",
    "Built real-time AI app for 10M+ concurrent users — Kumbh Mela 2025",
    "Skills: Python  •  TensorFlow  •  SQL  •  Tableau  •  Power BI  •  Flutter",
    "Team Lead — Software Developer at Vershama Tech Pvt. Ltd.",
]
for i, cred in enumerate(creds):
    y = 1.4 + i * 0.9
    add_rect(slide, 0.5, y + 0.15, 0.35, 0.35, GOOGLE_BLUE)
    add_text(slide, cred, 1.1, y, 11.5, 0.6,
             size=17, color=WHITE)

add_rect(slide, 0.5, 6.2, 12.3, 0.9, RGBColor(0x2a, 0x2a, 0x2a))
add_text(slide,
         "I don't just design AI systems. I deploy them at scale.",
         0.7, 6.3, 12, 0.65, size=20, bold=True, italic=True,
         color=GOOGLE_GREEN, align=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 9 — ASK
# ============================================================
slide = add_slide(prs)
set_bg(slide, DARK)
add_rect(slide, 0, 0, 13.33, 0.08, GOOGLE_RED)

add_text(slide, "What I Am Asking", 0.5, 0.3, 12, 0.8,
         size=36, bold=True, color=WHITE)

options = [
    (GOOGLE_BLUE,  "Option 1",
     "Integrate EduLens as a Google Workspace for Education add-on"),
    (GOOGLE_GREEN, "Option 2",
     "Deploy as a Gemini-powered analytics layer inside Google Classroom"),
    (GOOGLE_RED,   "Option 3",
     "Pilot with 1 institution — I will lead the full build and deployment"),
]
for i, (color, title, desc) in enumerate(options):
    y = 1.6 + i * 1.5
    add_rect(slide, 0.5, y, 12.3, 1.1, RGBColor(0x2a, 0x2a, 0x2a))
    add_text(slide, title, 0.8, y + 0.1, 2.5, 0.5,
             size=18, bold=True, color=color)
    add_text(slide, desc, 3.2, y + 0.1, 9.3, 0.9,
             size=18, color=WHITE)

add_rect(slide, 0.5, 6.3, 12.3, 0.8, GOOGLE_BLUE)
add_text(slide, "Give me the resources. I will give you the product.",
         0.7, 6.38, 12, 0.6, size=20, bold=True, italic=True,
         color=WHITE, align=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 10 — CLOSE
# ============================================================
slide = add_slide(prs)
set_bg(slide, DARK)
add_rect(slide, 0, 0, 13.33, 0.15, GOOGLE_BLUE)
add_rect(slide, 0, 7.35, 13.33, 0.15, GOOGLE_BLUE)

add_text(slide, "EduLens AI", 1, 0.8, 11, 1.2,
         size=52, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_text(slide, '"Google Classroom shows who attended.',
         1, 2.2, 11, 0.6, size=22, italic=True,
         color=RGBColor(0xBB, 0xBB, 0xBB), align=PP_ALIGN.CENTER)
add_text(slide, 'EduLens shows who understood —',
         1, 2.75, 11, 0.6, size=22, italic=True,
         color=WHITE, align=PP_ALIGN.CENTER)
add_text(slide, 'and who needs help before it\'s too late."',
         1, 3.3, 11, 0.6, size=22, italic=True,
         color=WHITE, align=PP_ALIGN.CENTER)

add_rect(slide, 4.0, 4.3, 5.33, 0.06, GOOGLE_BLUE)

add_text(slide, "Vaishnavi Singasani",
         1, 4.7, 11, 0.6, size=22, bold=True,
         color=WHITE, align=PP_ALIGN.CENTER)
add_text(slide, "vaishnaviv0809@gmail.com",
         1, 5.2, 11, 0.5, size=16,
         color=GOOGLE_BLUE, align=PP_ALIGN.CENTER)
add_text(slide, "linkedin.com/in/singasani-vaishnavi0031",
         1, 5.65, 11, 0.5, size=16,
         color=GOOGLE_BLUE, align=PP_ALIGN.CENTER)

# Save
output = r'C:\Users\vaish\OneDrive\Documents\Ai project\EduLens_AI_Pitch_Deck.pptx'
prs.save(output)
print("Pitch deck saved!")
print(f"Location: {output}")
print("10 slides complete — ready for Google!")
